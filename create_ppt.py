from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_COLOR = RGBColor(0xE8, 0xB4, 0x4B)  # Gold
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG = RGBColor(0x16, 0x21, 0x3E)
BLUE_ACCENT = RGBColor(0x0F, 0x3D, 0x6E)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_bullet_points(slide, left, top, width, height, items, font_size=16, color=WHITE, icon_color=ACCENT_COLOR):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        # Bullet icon
        run1 = p.add_run()
        run1.text = "  " + chr(0x2022) + "  "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = icon_color
        run1.font.bold = True
        # Text
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
    return txBox


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_COLOR
    shape.line.fill.background()
    return shape


def add_slide_number(slide, num, total):
    add_text_box(slide, Inches(12.3), Inches(7.0), Inches(1), Inches(0.4),
                 f"{num}/{total}", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)


TOTAL_SLIDES = 12

# ============================================================
# SLIDE 1 - TITLE SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Decorative top bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

# Title
add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "URBANTIME", font_size=56, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(2.6), Inches(2.3))

add_text_box(slide, Inches(1.5), Inches(2.9), Inches(10), Inches(1),
             "Luxury Watch E-Commerce Platform", font_size=30, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
             "A Full-Stack Web Application for Premium Timepiece Retail", font_size=18, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Info cards at bottom
for i, (label, value) in enumerate([
    ("Project Type", "Full-Stack Web Application"),
    ("Domain", "E-Commerce"),
    ("Team", "Final Year Project")
]):
    x = Inches(2.5 + i * 3)
    card = add_shape(slide, x, Inches(5.0), Inches(2.5), Inches(1.4), CARD_BG)
    add_text_box(slide, x, Inches(5.1), Inches(2.5), Inches(0.5),
                 label, font_size=13, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(5.5), Inches(2.5), Inches(0.6),
                 value, font_size=15, color=WHITE, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 1, TOTAL_SLIDES)

# ============================================================
# SLIDE 2 - INTRODUCTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "INTRODUCTION", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

# Left content
intro_text = (
    "UrbanTime is a full-stack luxury watch e-commerce platform designed to provide "
    "a seamless online shopping experience for premium timepieces. The platform caters "
    "to watch enthusiasts looking for curated collections from renowned brands worldwide."
)
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.5),
             intro_text, font_size=16, color=LIGHT_GRAY)

add_bullet_points(slide, Inches(0.8), Inches(3.2), Inches(5.5), Inches(3.5), [
    "Online luxury watch marketplace with curated collections",
    "Separate customer storefront and admin management panel",
    "Secure user authentication with JWT-based authorization",
    "Integrated Stripe payment gateway for online transactions",
    "AI-powered chatbot for real-time customer assistance",
    "Responsive design with dark/light theme support",
    "MongoDB-backed scalable database architecture",
], font_size=15)

# Right side - Key highlights cards
highlights = [
    ("Full-Stack", "React + Node.js + MongoDB"),
    ("Secure", "JWT + Bcrypt + Stripe"),
    ("Smart", "AI Chatbot Assistant"),
    ("Modern", "Tailwind CSS + Vite"),
]
for i, (title, desc) in enumerate(highlights):
    y = Inches(1.5 + i * 1.35)
    card = add_shape(slide, Inches(7.2), y, Inches(5.2), Inches(1.15), CARD_BG)
    add_text_box(slide, Inches(7.5), y + Inches(0.1), Inches(4.6), Inches(0.5),
                 title, font_size=20, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, Inches(7.5), y + Inches(0.55), Inches(4.6), Inches(0.5),
                 desc, font_size=14, color=LIGHT_GRAY)

add_slide_number(slide, 2, TOTAL_SLIDES)

# ============================================================
# SLIDE 3 - OBJECTIVES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "OBJECTIVES", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

objectives = [
    ("01", "Build a Scalable E-Commerce Platform",
     "Develop a full-stack web application using MERN stack for selling luxury watches online with a clean, modern UI."),
    ("02", "Implement Secure Authentication",
     "Create a robust JWT-based authentication system with password encryption using Bcrypt to protect user data."),
    ("03", "Enable Online Payment Processing",
     "Integrate Stripe payment gateway to support secure online transactions with payment confirmation workflow."),
    ("04", "Develop an Admin Management Panel",
     "Build a separate admin dashboard for managing watch inventory, viewing orders, and handling bookings."),
    ("05", "Provide AI-Powered Customer Support",
     "Implement an intelligent chatbot to assist users with product discovery, order tracking, and FAQs."),
    ("06", "Ensure Responsive & Accessible Design",
     "Create a mobile-first responsive interface with dark/light theme support for enhanced user experience."),
]

for i, (num, title, desc) in enumerate(objectives):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.6 + row * 2.8)

    card = add_shape(slide, x, y, Inches(3.8), Inches(2.4), CARD_BG)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_COLOR
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.9), y + Inches(0.2), Inches(2.7), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.4), Inches(1.3),
                 desc, font_size=13, color=LIGHT_GRAY)

add_slide_number(slide, 3, TOTAL_SLIDES)

# ============================================================
# SLIDE 4 - KEY FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "KEY FEATURES", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

features = [
    ("User Authentication", "Secure registration & login with JWT tokens and Bcrypt password hashing"),
    ("Watch Catalog", "Browse watches by category (Men/Women) and by brand with filtering & pagination"),
    ("Shopping Cart", "Add, update, remove items with real-time quantity management and price calculation"),
    ("Stripe Payments", "Secure online payment processing with Stripe Checkout sessions and payment confirmation"),
    ("Admin Panel", "Add/delete watches, manage product listings, view and process customer orders"),
    ("AI Chatbot", "Interactive chatbot for product recommendations, order tracking, and customer support"),
    ("Dark/Light Theme", "Toggle between dark and light modes with persistent theme preference"),
    ("Responsive Design", "Mobile-first layout with Tailwind CSS, hamburger menu for mobile navigation"),
    ("Order Management", "Complete order lifecycle with status tracking (Pending, Processing, Shipped, Delivered)"),
    ("Image Upload", "Multer-based image upload for watch product images with static file serving"),
]

for i, (title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.5 + row * 1.1)

    card = add_shape(slide, x, y, Inches(5.9), Inches(0.95), CARD_BG)
    # Accent dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.32), Inches(0.15), Inches(0.15))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_COLOR
    dot.line.fill.background()

    add_text_box(slide, x + Inches(0.5), y + Inches(0.05), Inches(5.2), Inches(0.4),
                 title, font_size=15, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.45), Inches(5.2), Inches(0.5),
                 desc, font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 4, TOTAL_SLIDES)

# ============================================================
# SLIDE 5 - SYSTEM ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "SYSTEM ARCHITECTURE", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

# Three-tier architecture boxes
tiers = [
    ("FRONTEND (Client)", [
        "React.js + Vite",
        "Tailwind CSS",
        "React Router DOM",
        "Axios HTTP Client",
        "Context API (State)",
    ], Inches(0.5), BLUE_ACCENT),
    ("BACKEND (Server)", [
        "Node.js + Express.js",
        "RESTful API Routes",
        "JWT Authentication",
        "Multer File Upload",
        "Stripe Integration",
    ], Inches(4.5), RGBColor(0x2C, 0x3E, 0x50)),
    ("DATABASE (Storage)", [
        "MongoDB Atlas",
        "Mongoose ODM",
        "User Collection",
        "Watch/Order/Cart",
        "Cloud-hosted (Atlas)",
    ], Inches(8.5), RGBColor(0x1A, 0x5C, 0x3A)),
]

for title, items, x, bg_color in tiers:
    card = add_shape(slide, x, Inches(1.6), Inches(3.8), Inches(3.5), bg_color)
    add_text_box(slide, x, Inches(1.7), Inches(3.8), Inches(0.5),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(1.2), Inches(2.25), Inches(1.4))
    add_bullet_points(slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(2.5),
                      items, font_size=13, color=WHITE, icon_color=ACCENT_COLOR)

# Arrows between boxes
for ax in [Inches(4.35), Inches(8.35)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, Inches(3.0), Inches(0.3), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_COLOR
    arrow.line.fill.background()

# External services section
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(12), Inches(0.5),
             "EXTERNAL SERVICES & INTEGRATIONS", font_size=18, color=ACCENT_COLOR, bold=True)

ext_services = [
    ("Stripe API", "Payment processing\nvia Checkout Sessions"),
    ("MongoDB Atlas", "Cloud database\nhosting service"),
    ("JWT", "Token-based\nauthentication"),
    ("Multer", "File upload\nmiddleware"),
]

for i, (name, desc) in enumerate(ext_services):
    x = Inches(0.8 + i * 3.15)
    card = add_shape(slide, x, Inches(5.9), Inches(2.8), Inches(1.2), CARD_BG)
    add_text_box(slide, x, Inches(5.95), Inches(2.8), Inches(0.4),
                 name, font_size=14, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.35), Inches(2.8), Inches(0.6),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 5, TOTAL_SLIDES)

# ============================================================
# SLIDE 6 - TECHNOLOGY STACK
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "TECHNOLOGY STACK", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

tech_categories = [
    ("Frontend", [
        ("React.js 19", "UI component library"),
        ("Vite 7", "Build tool & dev server"),
        ("Tailwind CSS 4", "Utility-first CSS framework"),
        ("React Router 7", "Client-side routing"),
        ("Axios", "HTTP client for API calls"),
        ("Lucide React", "Icon components"),
        ("React Toastify", "Toast notifications"),
    ]),
    ("Backend", [
        ("Node.js", "JavaScript runtime"),
        ("Express.js 5", "Web application framework"),
        ("Mongoose 9", "MongoDB ODM"),
        ("JSON Web Token", "Authentication tokens"),
        ("Bcryptjs", "Password hashing"),
        ("Multer", "File upload handling"),
        ("Stripe", "Payment processing"),
    ]),
    ("Database & Tools", [
        ("MongoDB Atlas", "Cloud NoSQL database"),
        ("Git & GitHub", "Version control"),
        ("Nodemon", "Auto-restart dev server"),
        ("Dotenv", "Environment config"),
        ("Validator.js", "Input validation"),
        ("CORS", "Cross-origin resource sharing"),
        ("UUID", "Unique ID generation"),
    ]),
]

for col, (category, techs) in enumerate(tech_categories):
    x = Inches(0.5 + col * 4.2)

    add_text_box(slide, x, Inches(1.5), Inches(3.8), Inches(0.5),
                 category, font_size=20, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_accent_line(slide, x + Inches(1.2), Inches(2.0), Inches(1.4))

    for i, (tech, desc) in enumerate(techs):
        y = Inches(2.3 + i * 0.68)
        card = add_shape(slide, x, y, Inches(3.8), Inches(0.58), CARD_BG)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.02), Inches(1.8), Inches(0.3),
                     tech, font_size=13, color=WHITE, bold=True)
        add_text_box(slide, x + Inches(2.0), y + Inches(0.02), Inches(1.7), Inches(0.3),
                     desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 6, TOTAL_SLIDES)

# ============================================================
# SLIDE 7 - WORKFLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "USER WORKFLOW", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

# Flow steps
steps = [
    ("1", "Register / Login", "User creates account\nor logs in with JWT\nauthentication"),
    ("2", "Browse Catalog", "Explore watches by\ncategory or brand\nwith filtering"),
    ("3", "Add to Cart", "Select watches,\nadjust quantities,\nreview selections"),
    ("4", "Checkout", "Enter shipping &\nbilling details,\nreview order summary"),
    ("5", "Payment", "Pay via Stripe or\nchoose Cash on\nDelivery option"),
    ("6", "Order Confirmed", "Receive confirmation,\ntrack order status\nthrough lifecycle"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.3 + i * 2.15)

    # Card
    card = add_shape(slide, x, Inches(1.8), Inches(1.9), Inches(3.0), CARD_BG)

    # Step number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.65), Inches(1.5), Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_COLOR
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_BG
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x, Inches(2.2), Inches(1.9), Inches(0.5),
                 title, font_size=15, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.1), Inches(2.8), Inches(1.7), Inches(1.8),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between steps
    if i < 5:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.0), Inches(3.0), Inches(0.18), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_COLOR
        arrow.line.fill.background()

# Admin workflow section
add_text_box(slide, Inches(0.8), Inches(5.2), Inches(12), Inches(0.5),
             "ADMIN WORKFLOW", font_size=20, color=ACCENT_COLOR, bold=True)

admin_steps = [
    ("Admin Login", "Access admin panel\nwith credentials"),
    ("Add Products", "Upload watches with\nimages, prices, brands"),
    ("Manage Orders", "View, update status,\nprocess bookings"),
    ("Inventory", "List/delete products,\nmonitor stock"),
]

for i, (title, desc) in enumerate(admin_steps):
    x = Inches(0.8 + i * 3.15)
    card = add_shape(slide, x, Inches(5.7), Inches(2.8), Inches(1.3), BLUE_ACCENT)
    add_text_box(slide, x, Inches(5.75), Inches(2.8), Inches(0.4),
                 title, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.2), Inches(2.8), Inches(0.6),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 7, TOTAL_SLIDES)

# ============================================================
# SLIDE 8 - PAYMENT & SECURITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "PAYMENT & SECURITY", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

# Payment section
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(5), Inches(0.5),
             "PAYMENT INTEGRATION", font_size=22, color=WHITE, bold=True)

payment_items = [
    "Stripe Checkout Sessions for secure online payments",
    "Support for card payments with real-time processing",
    "Cash on Delivery (COD) as alternative payment method",
    "Payment confirmation via Stripe session verification",
    "Order status updates based on payment status (Paid/Unpaid)",
    "Automatic tax calculation (8% tax rate) on orders",
    "Shipping charge logic - free shipping above threshold",
    "Success and cancel URL redirects after payment",
]
add_bullet_points(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(4.5),
                  payment_items, font_size=14, icon_color=GREEN)

# Security section
add_text_box(slide, Inches(7.2), Inches(1.5), Inches(5), Inches(0.5),
             "SECURITY MEASURES", font_size=22, color=WHITE, bold=True)

security_items = [
    "JWT (JSON Web Token) based authentication system",
    "Bcrypt password hashing with salt rounds (10)",
    "Auth middleware protecting sensitive API routes",
    "Token expiry set to 24 hours for session security",
    "Protected frontend routes redirect unauthenticated users",
    "CORS enabled for controlled cross-origin access",
    "Input validation using Validator.js library",
    "Sensitive credentials stored in environment variables",
]
add_bullet_points(slide, Inches(7.2), Inches(2.1), Inches(5.5), Inches(4.5),
                  security_items, font_size=14, icon_color=RED)

# Bottom info bar
card = add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), CARD_BG)
add_text_box(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7),
             "Stripe processes payments on its PCI-compliant servers — no sensitive card data touches our backend. "
             "JWT tokens are verified on every protected API request, and passwords are never stored in plain text.",
             font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 8, TOTAL_SLIDES)

# ============================================================
# SLIDE 9 - DATABASE DESIGN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
             "DATABASE SCHEMA", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

collections = [
    ("User Collection", [
        "_id (ObjectId)",
        "name (String, required)",
        "email (String, unique)",
        "password (String, hashed)",
        "timestamps (createdAt, updatedAt)",
    ]),
    ("Watch Collection", [
        "_id (ObjectId)",
        "name, description (String)",
        "price (Number, required)",
        "category (men/women/brand)",
        "brandName (String, optional)",
        "image (String, URL path)",
    ]),
    ("Order Collection", [
        "orderId (String, unique)",
        "user (ObjectId ref -> User)",
        "name, email, phone, address",
        "items[] (product array)",
        "totalAmount, taxAmount, finalAmount",
        "paymentMethod, paymentStatus",
        "orderStatus (Pending/Confirmed/etc)",
    ]),
    ("Cart Collection", [
        "user (ObjectId ref -> User, unique)",
        "items[] (productId, name, img,",
        "  price, qty, description)",
        "timestamps (createdAt, updatedAt)",
    ]),
]

for i, (title, fields) in enumerate(collections):
    x = Inches(0.5 + (i % 2) * 6.3)
    y = Inches(1.5 + (i // 2) * 2.9)

    card = add_shape(slide, x, y, Inches(6.0), Inches(2.6), CARD_BG)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.6), Inches(0.4),
                 title, font_size=18, color=ACCENT_COLOR, bold=True)
    add_accent_line(slide, x + Inches(0.2), y + Inches(0.55), Inches(2))
    add_bullet_points(slide, x + Inches(0.2), y + Inches(0.7), Inches(5.6), Inches(1.8),
                      fields, font_size=12, color=LIGHT_GRAY, icon_color=ACCENT_COLOR)

add_slide_number(slide, 9, TOTAL_SLIDES)

# ============================================================
# SLIDE 10 - ADVANTAGES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "ADVANTAGES", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

advantages = [
    ("Scalable Architecture", "MERN stack enables easy horizontal scaling and modular development with separate frontend, admin, and backend services."),
    ("Real-Time Cart Management", "Context API provides instant cart updates across all components without page reloads, enhancing UX."),
    ("AI-Powered Assistance", "Integrated chatbot reduces support load by handling product queries, recommendations, and order tracking automatically."),
    ("Secure Payment Flow", "Stripe handles all sensitive card data on PCI-compliant servers, minimizing security liability for the application."),
    ("Separate Admin Panel", "Dedicated admin interface allows inventory and order management without exposing admin functions to customers."),
    ("Modern Developer Experience", "Vite provides lightning-fast HMR builds, and Tailwind CSS enables rapid, consistent UI development."),
    ("Mobile-First Design", "Responsive layout with hamburger navigation ensures seamless experience across all device sizes."),
    ("Dual Payment Options", "Supporting both online (Stripe) and Cash on Delivery provides flexibility for diverse customer preferences."),
]

for i, (title, desc) in enumerate(advantages):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.5 + row * 1.4)

    card = add_shape(slide, x, y, Inches(5.9), Inches(1.2), CARD_BG)
    # Checkmark
    add_text_box(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.4), Inches(0.4),
                 chr(0x2713), font_size=20, color=GREEN, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.05), Inches(5.2), Inches(0.35),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.45), Inches(5.2), Inches(0.7),
                 desc, font_size=12, color=LIGHT_GRAY)

add_slide_number(slide, 10, TOTAL_SLIDES)

# ============================================================
# SLIDE 11 - FUTURE SCOPE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "FUTURE SCOPE", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

add_text_box(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.5),
             "Planned enhancements and potential expansions for the UrbanTime platform:",
             font_size=16, color=LIGHT_GRAY)

future_items = [
    ("AI/ML Recommendations", "Implement collaborative filtering and content-based recommendation engine using user browsing and purchase history for personalized watch suggestions."),
    ("Real-Time Notifications", "Add WebSocket-based push notifications for order status updates, new arrivals, price drops, and promotional offers."),
    ("Multi-Vendor Marketplace", "Expand to support multiple sellers/brands with vendor dashboards, commission management, and independent inventory control."),
    ("Wishlist & Reviews", "Allow users to save watches to wishlists and submit authenticated product reviews with star ratings and images."),
    ("Advanced Analytics", "Integrate analytics dashboards for admin with sales trends, revenue charts, customer demographics, and inventory forecasting."),
    ("Mobile Application", "Develop native iOS and Android apps using React Native, sharing business logic with the existing web platform."),
    ("Multi-Currency & i18n", "Support international payments with currency conversion, multi-language UI, and region-specific shipping rules."),
    ("Subscription & Loyalty", "Introduce membership tiers, loyalty points, referral programs, and recurring subscription boxes for collectors."),
]

for i, (title, desc) in enumerate(future_items):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(2.0 + row * 1.35)

    card = add_shape(slide, x, y, Inches(5.9), Inches(1.15), CARD_BG)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.05), Inches(0.4), Inches(0.35),
                 chr(0x2192), font_size=18, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.05), Inches(5.2), Inches(0.35),
                 title, font_size=15, color=ACCENT_COLOR, bold=True)
    add_text_box(slide, x + Inches(0.5), y + Inches(0.42), Inches(5.2), Inches(0.7),
                 desc, font_size=11, color=LIGHT_GRAY)

add_slide_number(slide, 11, TOTAL_SLIDES)

# ============================================================
# SLIDE 12 - CONCLUSION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Pt(6), ACCENT_COLOR)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.7),
             "CONCLUSION", font_size=36, color=ACCENT_COLOR, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.1), Inches(2))

conclusion_text = (
    "The UrbanTime luxury watch e-commerce platform successfully demonstrates a modern, "
    "full-stack web application built with the MERN stack. The project integrates all essential "
    "e-commerce functionalities including user authentication, product catalog management, "
    "shopping cart, secure Stripe payment processing, and an intelligent chatbot assistant."
)
add_text_box(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(1.2),
             conclusion_text, font_size=17, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Summary stats
stats = [
    ("3", "Applications\nFrontend + Admin + Backend"),
    ("4", "Collections\nUser, Watch, Order, Cart"),
    ("15+", "API Endpoints\nRESTful architecture"),
    ("10", "Key Features\nAuth, Cart, Pay, Chatbot"),
]

for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.15)
    card = add_shape(slide, x, Inches(3.0), Inches(2.8), Inches(1.5), CARD_BG)
    add_text_box(slide, x, Inches(3.05), Inches(2.8), Inches(0.6),
                 num, font_size=36, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.65), Inches(2.8), Inches(0.7),
                 label, font_size=13, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Key takeaways
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(0.5),
             "KEY TAKEAWAYS", font_size=20, color=ACCENT_COLOR, bold=True)

takeaways = [
    "Successfully built a production-ready e-commerce platform using modern web technologies",
    "Implemented secure authentication and payment systems following industry best practices",
    "Designed a scalable architecture that separates concerns across frontend, admin, and backend",
    "Created an AI-powered chatbot to enhance customer engagement and reduce support overhead",
]
add_bullet_points(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(2),
                  takeaways, font_size=15, icon_color=GREEN)

# Thank you
add_text_box(slide, Inches(1.5), Inches(6.8), Inches(10.3), Inches(0.5),
             "Thank You", font_size=28, color=ACCENT_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)

# Save
output_path = r"C:\Users\Suraj Singh\Desktop\E-commerce-Website\UrbanTime_FinalYearProject.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
